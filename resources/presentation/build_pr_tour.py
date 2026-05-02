"""Build a PR-tour slide deck matching the SecuriThings workshop deck style.

Each slide is identical to the existing repo-tour slides (slides 5-22 in
GitHub-Workshop-SecuriThings.pptx):
  - dark #0D1117 background
  - full-bleed top image (PR page screenshot, 10" x 4.81")
  - 👆 emoji (Segoe UI Emoji, 44pt, soft shadow) anchored under the target
  - bottom caption strip: orange section name + white description

Output: github-pr-tour.pptx (10 slides covering tabs + sidebar sections)
"""
from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

HERE = Path(__file__).parent
IMG = HERE / "github-pr-ui.png"
OUT = HERE / "github-pr-tour.pptx"

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

# Slide dimensions (matches reference deck): 10" x 5.625" = 9144000 x 5143500 EMU
SLIDE_W = 9144000
SLIDE_H = 5143500

# Image area: full width, 4.81" tall (same as reference)
IMG_W = 9144000
IMG_H = 4398264

# Emoji: 0.6" x 0.6" (slightly smaller than the reference deck's 0.7"
# because the PR sidebar items are tightly packed; bigger emoji bodies
# would visually overlap the neighbouring section)
EMOJI_W = 548640
EMOJI_H = 548640
EMOJI_FONT_PT = 38

# Source image dimensions in pixels (the cropped PR screenshot)
SRC_PX_W = 1920
SRC_PX_H = 920

# EMU per pixel in the placed image
PX2EMU_X = IMG_W / SRC_PX_W   # 4762.5
PX2EMU_Y = IMG_H / SRC_PX_H   # 4780.7


def emoji_pos_below(target_cx_px: int, target_bottom_y_px: int) -> tuple[int, int]:
    """Position emoji centered horizontally under target — for top-of-page tabs."""
    x = int(target_cx_px * PX2EMU_X - EMOJI_W / 2)
    y = int(target_bottom_y_px * PX2EMU_Y + 45720)
    return x, y


def emoji_pos_left(target_left_px: int, target_center_y_px: int) -> tuple[int, int]:
    """Position emoji to the LEFT of the target, vertically centred on it.

    Used for sidebar headings: a 👉 finger floats in the empty gutter between
    the conversation column and the sidebar, pointing right at the heading.
    """
    gap_emu = 45720  # ~0.05" between fingertip and target text
    x = int(target_left_px * PX2EMU_X - EMOJI_W - gap_emu)
    y = int(target_center_y_px * PX2EMU_Y - EMOJI_H / 2)
    return x, y


# (name, description, emoji_char, anchor_kind, target coords)
#   anchor_kind="below"  → (target_center_x_px, target_bottom_y_px)
#   anchor_kind="left"   → (target_left_x_px,   target_center_y_px)
TOUR = [
    ("Conversation",
     "The unified PR thread — description, comments, reviews, and event log.",
     "👆", "below", 420, 348),
    ("Commits",
     "Browse the individual commits in this PR. Useful for stacked-diff reviews.",
     "👆", "below", 567, 348),
    ("Checks",
     "Status of every CI run for this PR: tests, linters, security scans, deploys.",
     "👆", "below", 696, 348),
    ("Files changed",
     "The full diff. Leave inline comments, suggested changes, and approvals here.",
     "👆", "below", 838, 348),
    ("Reviewers",
     "Request review from people or teams. CODEOWNERS auto-fills here.",
     "👉", "left", 1247, 383),
    ("Assignees",
     "Who is responsible for landing this PR. Often the author, sometimes a shepherd.",
     "👉", "left", 1247, 460),
    ("Labels",
     "Categorise the PR: bug, feature, area/api, needs-triage. Drives filters & automation.",
     "👉", "left", 1247, 536),
    ("Projects",
     "Add to a GitHub Project board — Kanban, roadmap, or sprint tracking.",
     "👉", "left", 1247, 613),
    ("Milestone",
     "Group PRs by release or iteration (e.g. v1.2.0, Sprint 14).",
     "👉", "left", 1247, 689),
    ("Development",
     "Link issues. \"Closes #123\" in the description auto-closes them on merge.",
     "👉", "left", 1247, 766),
]


def add_background(slide, hex_color: str) -> None:
    """Set solid slide background color via cSld/bg."""
    cSld = slide._element.find(qn("p:cSld"))
    bg_xml = (
        f'<p:bg xmlns:p="{P_NS}" xmlns:a="{A_NS}">'
        f'  <p:bgPr>'
        f'    <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
        f'    <a:effectLst/>'
        f'  </p:bgPr>'
        f'</p:bg>'
    )
    bg = etree.fromstring(bg_xml)
    cSld.insert(0, bg)


def add_emoji(slide, x: int, y: int, glyph: str = "👆") -> None:
    """Add an emoji textbox with soft drop-shadow at the given EMU position."""
    tx = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(EMOJI_W), Emu(EMOJI_H))
    tf = tx.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    from pptx.enum.text import PP_ALIGN
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = glyph
    run.font.name = "Segoe UI Emoji"
    run.font.size = Pt(EMOJI_FONT_PT)
    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    bodyPr.set("anchor", "ctr")

    spPr = tx._element.find(qn("p:spPr"))
    for el in spPr.findall(qn("a:effectLst")):
        spPr.remove(el)
    effect_xml = (
        f'<a:effectLst xmlns:a="{A_NS}">'
        f'  <a:outerShdw blurRad="76200" dist="25400" dir="5400000" algn="bl" rotWithShape="0">'
        f'    <a:srgbClr val="000000"><a:alpha val="70000"/></a:srgbClr>'
        f'  </a:outerShdw>'
        f'</a:effectLst>'
    )
    spPr.append(etree.fromstring(effect_xml))


def add_caption(slide, name: str, desc: str) -> None:
    """Bottom caption strip: orange name + gray dash + white description."""
    # Position matches reference: x=365760, y=4443984, cx=8412480, cy=653796
    tx = slide.shapes.add_textbox(
        Emu(365760), Emu(4443984), Emu(8412480), Emu(653796)
    )
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    bodyPr.set("anchor", "ctr")

    p = tf.paragraphs[0]

    r1 = p.add_run()
    r1.text = name
    r1.font.name = "Calibri"
    r1.font.size = Pt(20)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(0xF0, 0x88, 0x3E)  # GitHub orange

    r2 = p.add_run()
    r2.text = "  —  "
    r2.font.name = "Calibri"
    r2.font.size = Pt(16)
    r2.font.color.rgb = RGBColor(0xC9, 0xD1, 0xD9)

    r3 = p.add_run()
    r3.text = desc
    r3.font.name = "Calibri"
    r3.font.size = Pt(14)
    r3.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def build_slide(prs, name: str, desc: str, glyph: str, kind: str, a: int, b: int):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    add_background(slide, "0D1117")

    slide.shapes.add_picture(
        str(IMG), Emu(0), Emu(0), width=Emu(IMG_W), height=Emu(IMG_H)
    )

    if kind == "below":
        ex, ey = emoji_pos_below(a, b)
    elif kind == "left":
        ex, ey = emoji_pos_left(a, b)
    else:
        raise ValueError(f"unknown anchor kind: {kind}")
    add_emoji(slide, ex, ey, glyph)

    add_caption(slide, name, desc)
    return slide


def main() -> None:
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    for name, desc, glyph, kind, a, b in TOUR:
        build_slide(prs, name, desc, glyph, kind, a, b)

    prs.save(OUT)
    print(f"Wrote {OUT} ({len(TOUR)} slides)")


if __name__ == "__main__":
    main()
