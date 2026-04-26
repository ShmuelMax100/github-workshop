"""
Build a vivid, modern PowerPoint deck from gitlab-to-github-slides.md.

Design language:
- Title slide: deep gradient + decorative orbs + display title
- Content slides: white canvas, colored accent rail, bullet cards,
  subtle corner blob, section color rotated per group
- Interlude slides: full-bleed two-color gradient + giant centered quote
  + scattered translucent orbs (used as humor breaks)

Usage:
    pip install python-pptx
    python build_pptx.py
"""

from __future__ import annotations

import re
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import nsmap, qn
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).parent
SRC = HERE / "gitlab-to-github-slides.md"
OUT = HERE / "gitlab-to-github.pptx"

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ---------- Palette ----------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x0F, 0x17, 0x2A)        # near-black ink
NAVY = RGBColor(0x0D, 0x1B, 0x2A)
SOFT = RGBColor(0xF5, 0xF7, 0xFA)
MUTED = RGBColor(0x6B, 0x72, 0x80)

# Section accent colors (hex strings for XML), rotated by section_id
SECTIONS = {
    "intro":      ("4F46E5", "312E81"),  # indigo
    "pr":         ("2563EB", "1E3A8A"),  # blue
    "cli":        ("0D9488", "134E4A"),  # teal
    "protect":    ("059669", "064E3B"),  # emerald
    "ui":         ("7C3AED", "4C1D95"),  # violet
    "wrap":       ("0EA5E9", "0C4A6E"),  # sky
}

# Interlude gradients (vibrant, full-bleed)
INTERLUDE_PALETTE = [
    ("F97316", "BE185D"),  # orange → pink
    ("8B5CF6", "1E1B4B"),  # violet → midnight
    ("F59E0B", "7C2D12"),  # amber → rust
    ("EC4899", "831843"),  # pink → wine
    ("06B6D4", "164E63"),  # cyan → deep teal
    ("10B981", "064E3B"),  # emerald → forest
]

# Map slide index (1-based) to section key for content slides
def section_for(i: int) -> str:
    if i == 1: return "intro"
    if 2 <= i <= 5: return "intro"
    if 7 <= i <= 13: return "pr"
    if 15 <= i <= 16: return "cli"
    if 18 <= i <= 22: return "protect"
    if 24 <= i <= 25: return "ui"
    if 27 <= i <= 30: return "wrap"
    return "intro"


# ---------- XML helpers ----------
def _spPr(shape):
    return shape._element.find(qn("p:spPr"))


def _clear_fills(spPr):
    for tag in ("solidFill", "gradFill", "blipFill", "noFill", "pattFill", "grpFill"):
        for el in spPr.findall(qn(f"a:{tag}")):
            spPr.remove(el)


def _insert_fill(spPr, fill_el):
    """Insert a fill element in the correct schema position (after geometry, before <a:ln>)."""
    # Find <a:ln> if present — fill must come before it.
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        ln.addprevious(fill_el)
        return
    # Else, append after geometry.
    spPr.append(fill_el)


def set_gradient_fill(shape, color_a: str, color_b: str, angle: int = 2700000):
    spPr = _spPr(shape)
    _clear_fills(spPr)
    xml = (
        f'<a:gradFill xmlns:a="{A_NS}" rotWithShape="1">'
        f'  <a:gsLst>'
        f'    <a:gs pos="0"><a:srgbClr val="{color_a}"/></a:gs>'
        f'    <a:gs pos="100000"><a:srgbClr val="{color_b}"/></a:gs>'
        f'  </a:gsLst>'
        f'  <a:lin ang="{angle}" scaled="0"/>'
        f'</a:gradFill>'
    )
    _insert_fill(spPr, etree.fromstring(xml))


def set_solid_alpha(shape, color_hex: str, alpha_pct: int):
    spPr = _spPr(shape)
    _clear_fills(spPr)
    xml = (
        f'<a:solidFill xmlns:a="{A_NS}">'
        f'  <a:srgbClr val="{color_hex}">'
        f'    <a:alpha val="{alpha_pct * 1000}"/>'
        f'  </a:srgbClr>'
        f'</a:solidFill>'
    )
    _insert_fill(spPr, etree.fromstring(xml))


def no_line(shape):
    shape.line.fill.background()


# ---------- Layout primitives ----------
def add_full_bleed_gradient(slide, prs, color_a, color_b, angle=2700000):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    no_line(rect)
    set_gradient_fill(rect, color_a, color_b, angle)
    return rect


def add_orb(slide, cx_in, cy_in, d_in, color_hex, alpha_pct):
    d = Inches(d_in)
    orb = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(cx_in - d_in / 2), Inches(cy_in - d_in / 2), d, d
    )
    no_line(orb)
    set_solid_alpha(orb, color_hex, alpha_pct)
    return orb


def add_text(slide, left, top, width, height, text, *,
             font="Segoe UI", size=18, bold=False, color=INK, align=None,
             italic=False, line_spacing=None):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tx


def add_rich_text(slide, left, top, width, height, segments, *,
                  size=18, color=INK, font="Segoe UI", align=None, line_spacing=None):
    """segments: list of (text, bold)"""
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    for text, bold in segments:
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tx


# ---------- Parsing ----------
def parse_slides(md: str) -> list[dict]:
    chunks = re.split(r"^## Slide \d+ — ", md, flags=re.MULTILINE)[1:]
    slides = []
    for chunk in chunks:
        lines = chunk.splitlines()
        title = lines[0].strip()
        body = "\n".join(lines[1:])

        is_interlude = title.startswith("[INTERLUDE]")
        if is_interlude:
            title = title.replace("[INTERLUDE]", "").strip()

        # Split body at speaker-notes marker so we never mix them with quote/bullets
        body_main, _, _ = body.partition("**Speaker notes:**")

        # Bullets (or, for interludes, **bold** lines as quote lines)
        bullets = [m.group(1).strip()
                   for m in re.finditer(r"^- (.+)$", body_main, flags=re.MULTILINE)]

        quote_lines = []
        attribution = ""
        if is_interlude:
            for m in re.finditer(r"^\*\*(.+?)\*\*\s*$", body_main, flags=re.MULTILINE):
                quote_lines.append(m.group(1).strip())
            am = re.search(r"^\*— (.+)\*\s*$", body_main, flags=re.MULTILINE)
            if am:
                attribution = am.group(1).strip()

        notes_match = re.search(
            r"\*\*Speaker notes:\*\*\s*\n(.+?)(?=\n\n\*\*|\n\n---|\Z)",
            body, flags=re.DOTALL,
        )
        notes = notes_match.group(1).strip() if notes_match else ""
        notes = re.sub(r"\s+\n", "\n", notes)
        notes = re.sub(r"\n+", " ", notes).strip()

        slides.append({
            "title": title,
            "bullets": bullets,
            "quote_lines": quote_lines,
            "attribution": attribution,
            "notes": notes,
            "interlude": is_interlude,
        })
    return slides


# ---------- Slide builders ----------
def build_title_slide(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bleed_gradient(s, prs, "1E1B4B", "0F172A", angle=2700000)

    # decorative orbs
    add_orb(s, 12.2, 1.2, 3.4, "8B5CF6", 18)
    add_orb(s, 11.6, 6.4, 2.6, "06B6D4", 16)
    add_orb(s, 1.0, 6.6, 2.0, "F472B6", 14)
    add_orb(s, 0.6, 0.8, 1.2, "F59E0B", 22)

    # accent stripe
    stripe = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.45), Inches(0.45), Inches(2.6)
    )
    no_line(stripe)
    set_gradient_fill(stripe, "F472B6", "8B5CF6", angle=5400000)

    # eyebrow label
    add_text(s, Inches(1.6), Inches(2.3), Inches(11), Inches(0.4),
             "WORKSHOP  ·  GITLAB → GITHUB",
             font="Segoe UI", size=13, bold=True, color=RGBColor(0xF8, 0xB4, 0xD9))

    # display title
    add_text(s, Inches(1.6), Inches(2.7), Inches(11.2), Inches(2.0),
             "From GitLab to GitHub",
             font="Segoe UI Black", size=64, bold=True, color=WHITE,
             line_spacing=1.0)

    # subtitle
    add_text(s, Inches(1.6), Inches(4.2), Inches(11.2), Inches(0.7),
             "Daily Workflows, PRs, and Repo Lifecycle",
             font="Segoe UI Light", size=28, italic=True,
             color=RGBColor(0xCB, 0xD5, 0xE1))

    # tag chips
    chips = ["Pull Requests", "gh CLI", "Branch Protection", "CODEOWNERS", "UI Flows"]
    x = Inches(1.6)
    y = Inches(5.3)
    for label in chips:
        chip = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.0), Inches(0.45)
        )
        no_line(chip)
        set_solid_alpha(chip, "FFFFFF", 14)
        chip.adjustments[0] = 0.5
        ctf = chip.text_frame
        ctf.margin_left = Inches(0.15)
        ctf.margin_right = Inches(0.15)
        ctf.margin_top = Inches(0.05)
        ctf.margin_bottom = Inches(0.05)
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        crun = cp.add_run()
        crun.text = label
        crun.font.name = "Segoe UI Semibold"
        crun.font.size = Pt(13)
        crun.font.color.rgb = WHITE
        x += Inches(2.15)

    if data["notes"]:
        s.notes_slide.notes_text_frame.text = data["notes"]
    return s


def build_content_slide(prs, data, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    sec_key = section_for(page)
    color_a, color_b = SECTIONS[sec_key]

    # white background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    no_line(bg)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    # left accent rail (gradient)
    rail = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.45), prs.slide_height
    )
    no_line(rail)
    set_gradient_fill(rail, color_a, color_b, angle=5400000)

    # corner decoration: large translucent blob bottom-right
    add_orb(s, 13.4, 7.6, 4.2, color_a, 10)
    add_orb(s, 12.6, 7.0, 1.6, color_b, 14)

    # small accent square top-right
    acc = s.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(12.7), Inches(0.45), Inches(0.35), Inches(0.35)
    )
    no_line(acc)
    acc.fill.solid()
    acc.fill.fore_color.rgb = RGBColor.from_string(color_a)

    # slide number (faint, big)
    add_text(s, Inches(0.75), Inches(0.35), Inches(2.0), Inches(0.7),
             f"{page:02d}",
             font="Segoe UI Black", size=42, bold=True,
             color=RGBColor.from_string(color_a))

    # section eyebrow
    add_text(s, Inches(2.0), Inches(0.55), Inches(8), Inches(0.4),
             sec_key.upper() + "  ·  GITLAB → GITHUB",
             font="Segoe UI", size=11, bold=True, color=MUTED)

    # title
    add_text(s, Inches(0.85), Inches(1.05), Inches(11.6), Inches(1.0),
             data["title"],
             font="Segoe UI Semibold", size=34, bold=True, color=NAVY,
             line_spacing=1.05)

    # underline accent
    underline = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(1.95), Inches(0.7), Inches(0.07)
    )
    no_line(underline)
    underline.fill.solid()
    underline.fill.fore_color.rgb = RGBColor.from_string(color_a)

    # bullet cards
    bullets = data["bullets"]
    n = len(bullets)
    top = Inches(2.35)
    available = Inches(7.5 - 2.35 - 0.6)  # space for cards before footer
    gap = Inches(0.12)
    card_h_emu = (available - gap * (n - 1)) // max(n, 1)
    card_h = Emu(int(card_h_emu))

    left = Inches(0.85)
    width = Inches(11.6)

    for i, b in enumerate(bullets):
        y = top + Emu(int(card_h_emu) * i + (gap * i if i else 0))
        # card
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y, width, card_h)
        no_line(card)
        set_solid_alpha(card, color_a, 8)  # very subtle tint
        card.adjustments[0] = 0.18

        # accent dot
        dot_d = Inches(0.22)
        dot = s.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left + Inches(0.35),
            y + Emu(int(card_h_emu) // 2) - Inches(0.11),
            dot_d, dot_d,
        )
        no_line(dot)
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor.from_string(color_a)

        # bullet text (bold spans rendered)
        segments = []
        for part in re.split(r"(\*\*[^*]+\*\*)", b):
            if not part:
                continue
            if part.startswith("**") and part.endswith("**"):
                segments.append((part[2:-2], True))
            else:
                segments.append((part, False))
        tx = s.shapes.add_textbox(
            left + Inches(0.85),
            y + Inches(0.05),
            width - Inches(1.0),
            card_h - Inches(0.1),
        )
        tf = tx.text_frame
        tf.word_wrap = True
        # vertical-center text approximately by setting top margin
        tf.margin_top = Inches(0.0)
        tf.margin_bottom = Inches(0.0)
        p = tf.paragraphs[0]
        # nudge font size if many bullets
        size = 19 if n <= 4 else 17
        for text, bold in segments:
            run = p.add_run()
            run.text = text
            run.font.name = "Segoe UI"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = INK

    # footer
    add_text(s, Inches(0.85), Inches(7.05),
             prs.slide_width - Inches(1.7), Inches(0.3),
             f"GitLab → GitHub   ·   {page} / {total}",
             font="Segoe UI", size=10, color=MUTED)

    if data["notes"]:
        s.notes_slide.notes_text_frame.text = data["notes"]
    return s


def build_interlude_slide(prs, data, page, total, palette_idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    color_a, color_b = INTERLUDE_PALETTE[palette_idx % len(INTERLUDE_PALETTE)]
    add_full_bleed_gradient(s, prs, color_a, color_b, angle=2700000)

    # scattered orbs
    add_orb(s, 1.5, 1.3, 2.6, "FFFFFF", 12)
    add_orb(s, 12.0, 1.5, 1.8, "FFFFFF", 10)
    add_orb(s, 11.0, 6.4, 3.4, "FFFFFF", 8)
    add_orb(s, 0.5, 6.8, 1.4, "FFFFFF", 14)
    add_orb(s, 6.7, 0.6, 1.0, "FFFFFF", 18)

    # eyebrow
    add_text(s, Inches(0.85), Inches(0.55), Inches(10), Inches(0.4),
             "INTERLUDE",
             font="Segoe UI", size=12, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))

    # quote text — center on slide
    quote_block = "\n".join(data["quote_lines"]) if data["quote_lines"] else data["title"]
    tx = s.shapes.add_textbox(
        Inches(1.2), Inches(2.4), Inches(11.0), Inches(3.2)
    )
    tf = tx.text_frame
    tf.word_wrap = True
    for i, line in enumerate(data["quote_lines"] or [data["title"]]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.15
        run = p.add_run()
        run.text = line
        run.font.name = "Georgia"
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.italic = False
        run.font.color.rgb = WHITE

    # attribution
    if data["attribution"]:
        add_text(s, Inches(1.2), Inches(5.7), Inches(11.0), Inches(0.5),
                 "— " + data["attribution"],
                 font="Segoe UI Light", size=18, italic=True,
                 color=RGBColor(0xFB, 0xCF, 0xE8), align=PP_ALIGN.CENTER)

    # caption / title under attribution
    add_text(s, Inches(1.2), Inches(6.4), Inches(11.0), Inches(0.4),
             data["title"].upper(),
             font="Segoe UI", size=11, bold=True,
             color=RGBColor(0xFF, 0xE4, 0xE6), align=PP_ALIGN.CENTER)

    # footer
    add_text(s, Inches(0.85), Inches(7.05),
             prs.slide_width - Inches(1.7), Inches(0.3),
             f"INTERLUDE   ·   {page} / {total}",
             font="Segoe UI", size=10,
             color=RGBColor(0xFE, 0xCD, 0xD3))

    if data["notes"]:
        s.notes_slide.notes_text_frame.text = data["notes"]
    return s


def main():
    md = SRC.read_text(encoding="utf-8")
    slides = parse_slides(md)
    if not slides:
        raise SystemExit("No slides parsed.")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total = len(slides)
    interlude_idx = 0
    for i, data in enumerate(slides, start=1):
        if i == 1:
            build_title_slide(prs, data)
        elif data["interlude"]:
            build_interlude_slide(prs, data, i, total, interlude_idx)
            interlude_idx += 1
        else:
            build_content_slide(prs, data, i, total)

    prs.save(OUT)
    print(f"Wrote {OUT} ({total} slides)")


if __name__ == "__main__":
    main()
