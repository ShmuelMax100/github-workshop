"""Validate the generated pptx: slide count, titles, bullets, notes, overflow."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

PPTX = Path(__file__).parent / "gitlab-to-github.pptx"
MD = Path(__file__).parent / "gitlab-to-github-slides.md"

prs = Presentation(PPTX)
print(f"File: {PPTX.name}  size={PPTX.stat().st_size} bytes")
print(f"Slide size: {prs.slide_width/914400:.2f} x {prs.slide_height/914400:.2f} in")
print(f"Slide count: {len(prs.slides)}")

issues = []
for i, slide in enumerate(prs.slides, 1):
    title_txt = ""
    bullet_lines = []
    for shp in slide.shapes:
        if not shp.has_text_frame:
            continue
        for p_idx, para in enumerate(shp.text_frame.paragraphs):
            txt = "".join(r.text for r in para.runs)
            if not txt.strip():
                continue
            if not title_txt and any(r.font.size and r.font.size.pt >= 30 for r in para.runs):
                title_txt = txt
            elif txt.startswith("•"):
                bullet_lines.append(txt)
    notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
    print(f"\n--- Slide {i} ---")
    print(f"  title  : {title_txt[:90]}")
    print(f"  bullets: {len(bullet_lines)}")
    for b in bullet_lines:
        bl = b.lstrip("• ").strip()
        marker = "  !" if len(bl) > 110 else "   "
        print(f"  {marker} - {bl[:120]}")
        if len(bl) > 130:
            issues.append(f"slide {i} bullet too long ({len(bl)} chars)")
    print(f"  notes  : {len(notes)} chars")
    if i > 1 and len(bullet_lines) == 0:
        issues.append(f"slide {i} has no bullets")
    if i > 1 and (len(bullet_lines) < 3 or len(bullet_lines) > 5):
        issues.append(f"slide {i} bullets={len(bullet_lines)} (want 3-5)")
    if not notes.strip():
        issues.append(f"slide {i} missing notes")

# Check overflow: any text shape extending past slide edges
SW, SH = prs.slide_width, prs.slide_height
for i, slide in enumerate(prs.slides, 1):
    for shp in slide.shapes:
        if shp.left is None: continue
        right = shp.left + (shp.width or 0)
        bottom = shp.top + (shp.height or 0)
        if right > SW + 100 or bottom > SH + 100 or shp.left < -100 or shp.top < -100:
            issues.append(f"slide {i} shape '{shp.shape_type}' off-canvas: L={shp.left} T={shp.top} R={right} B={bottom}")

print("\n=== ISSUES ===")
if issues:
    for x in issues: print(" -", x)
else:
    print("none")
